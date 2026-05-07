"""Streamlit app for the Cloud.AI V2 agent-style FP&A workflow."""

from __future__ import annotations

import hashlib
import json
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from agent_controller import run_agent_workflow
from generate_summary import money, number, percent


st.set_page_config(page_title="Cloud.AI Board Reporting Agent", layout="wide")


def _json_text(data: dict[str, Any]) -> str:
    """Return pretty JSON for display and download."""

    return json.dumps(data, indent=2)


def _file_signature(uploaded_file) -> str:
    """Create a stable signature so Streamlit does not rerun analysis unnecessarily."""

    payload = uploaded_file.getbuffer()
    return hashlib.sha256(payload).hexdigest()


def _metric(agent_results: dict[str, Any], key: str) -> Any:
    """Return a summary metric from agent results."""

    return agent_results.get("metrics", {}).get("summary_metrics", {}).get(key)


def _status_badge(status: str) -> None:
    """Render a Streamlit status callout."""

    if status in {"Ready", "Green", "Ready for Stakeholders"}:
        st.success(status)
    elif status in {"Needs Review", "Amber"}:
        st.warning(status)
    else:
        st.error(status)


def _render_bullets(items: list[Any]) -> None:
    """Render a clean bullet list for strings or small dictionaries."""

    if not items:
        st.caption("No items flagged.")
        return
    for item in items:
        if isinstance(item, dict):
            st.write(f"- {json.dumps(item, default=str)}")
        else:
            st.write(f"- {item}")


def _build_pdf(report_markdown: str, agent_results: dict[str, Any]) -> bytes:
    """Create a compact board-report PDF from Markdown and agent review results."""

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=42, leftMargin=42, topMargin=42, bottomMargin=42)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("Cloud.AI Board Reporting Pack", styles["Title"]))
    story.append(Spacer(1, 12))

    summary_rows = [
        ["Readiness", str(agent_results.get("readiness", {}).get("readiness_score", "n/a"))],
        ["Model Health", str(agent_results.get("model_audit", {}).get("model_health_score", "n/a"))],
        ["Quality", str(agent_results.get("quality_review", {}).get("quality_score", "n/a"))],
        ["Ending ARR", money(_metric(agent_results, "ending_arr"))],
        ["Ending Cash", money(_metric(agent_results, "ending_cash"))],
        ["Funding Need", money(_metric(agent_results, "funding_need"))],
    ]
    table = Table(summary_rows, colWidths=[150, 300])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#E2E8F0")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#111827")),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("PADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 14))

    for raw_line in report_markdown.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("|") or line.startswith("---"):
            continue
        if line.startswith("# "):
            story.append(Paragraph(escape(line[2:]), styles["Title"]))
        elif line.startswith("## "):
            story.append(Spacer(1, 8))
            story.append(Paragraph(escape(line[3:]), styles["Heading2"]))
        elif line.startswith("- "):
            story.append(Paragraph(f"&#8226; {escape(line[2:])}", styles["BodyText"]))
        else:
            story.append(Paragraph(escape(line.replace("`", "")), styles["BodyText"]))

    doc.build(story)
    buffer.seek(0)
    return buffer.read()


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    """Add a simple PowerPoint bullet slide."""

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def _build_pptx(agent_results: dict[str, Any]) -> bytes:
    """Create a board-ready PPTX summary from agent outputs."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Cloud.AI Board Reporting Agent"
    title_slide.placeholders[1].text = "V2 Agent FP&A Analysis"

    _add_bullet_slide(
        prs,
        "Readiness and Model Health",
        [
            f"Readiness: {agent_results.get('readiness', {}).get('status')} ({agent_results.get('readiness', {}).get('readiness_score')}/100)",
            f"Model health: {agent_results.get('model_audit', {}).get('status')} ({agent_results.get('model_audit', {}).get('model_health_score')}/100)",
            f"Quality review: {agent_results.get('quality_review', {}).get('status')} ({agent_results.get('quality_review', {}).get('quality_score')}/100)",
        ],
    )

    _add_bullet_slide(
        prs,
        "Executive KPI Snapshot",
        [
            f"Ending ARR: {money(_metric(agent_results, 'ending_arr'))}",
            f"Ending MRR: {money(_metric(agent_results, 'ending_mrr'))}",
            f"Gross margin: {percent(_metric(agent_results, 'gross_margin'))}",
            f"Funding need: {money(_metric(agent_results, 'funding_need'))}",
        ],
    )

    finance = agent_results.get("finance_reasoning", {})
    executive = finance.get("executive_summary", {})
    _add_bullet_slide(
        prs,
        "Agent Analysis",
        executive.get("Facts", []) + executive.get("Interpretation", []) + executive.get("Recommendations", []),
    )

    _add_bullet_slide(
        prs,
        "Devil's Advocate",
        agent_results.get("devils_advocate", {}).get("questions_management_must_answer", [])[:5],
    )

    _add_bullet_slide(
        prs,
        "Recommended Actions",
        finance.get("recommended_actions", []),
    )

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _run_workflow_for_upload(uploaded_file) -> dict[str, Any]:
    """Persist the upload and run the V2 workflow."""

    work_dir = Path("streamlit_outputs")
    input_dir = work_dir / "input"
    output_dir = work_dir / "output"
    charts_dir = work_dir / "charts"
    reports_dir = work_dir / "reports"

    input_dir.mkdir(parents=True, exist_ok=True)
    workbook_path = input_dir / uploaded_file.name
    workbook_path.write_bytes(uploaded_file.getbuffer())

    return run_agent_workflow(
        workbook_path=workbook_path,
        output_dir=output_dir,
        reports_dir=reports_dir,
        charts_dir=charts_dir,
    )


st.title("Cloud.AI Board Reporting Agent")
st.caption("V2 agent-style SaaS FP&A workflow: readiness, model audit, KPI extraction, analysis, challenge review, and downloads")

tabs = st.tabs(
    [
        "Upload & Readiness",
        "Model Health Check",
        "KPI Dashboard",
        "Agent Analysis",
        "Devil's Advocate",
        "Downloads",
    ]
)

with tabs[0]:
    uploaded_file = st.file_uploader("Upload SaaS financial model workbook", type=["xlsx", "xlsm"])
    if uploaded_file:
        signature = _file_signature(uploaded_file)
        should_run = st.session_state.get("uploaded_signature") != signature

        if should_run:
            with st.spinner("Running V2 FP&A agent workflow..."):
                st.session_state["agent_results"] = _run_workflow_for_upload(uploaded_file)
                st.session_state["uploaded_signature"] = signature

        agent_results = st.session_state["agent_results"]
        readiness = agent_results.get("readiness", {})
        st.subheader("Workbook Readiness")
        col1, col2 = st.columns([1, 3])
        col1.metric("Readiness Score", f"{readiness.get('readiness_score', 0)}/100")
        with col2:
            _status_badge(readiness.get("status", "Not Ready"))

        st.write("Available sheets")
        st.write(", ".join(readiness.get("available_sheet_names", [])))
        st.write("Missing items")
        _render_bullets(readiness.get("missing_items", []))
        st.write("Manual review warnings")
        _render_bullets(readiness.get("manual_review_warnings", []))
    else:
        st.info("Upload the SaaS financial model workbook to generate the V2 board reporting output.")

agent_results = st.session_state.get("agent_results")

with tabs[1]:
    if not agent_results:
        st.info("Upload a workbook first.")
    else:
        audit = agent_results.get("model_audit", {})
        st.subheader("Model Health Check")
        col1, col2 = st.columns([1, 3])
        col1.metric("Health Score", f"{audit.get('model_health_score', 0)}/100")
        with col2:
            _status_badge(audit.get("status", "Red"))

        st.write("Issues")
        for issue in audit.get("issue_list", []):
            st.warning(f"{issue.get('severity')}: {issue.get('issue')}")
            st.json(issue)

        st.write("Recommended fixes")
        _render_bullets(audit.get("recommended_fixes", []))

with tabs[2]:
    if not agent_results:
        st.info("Upload a workbook first.")
    else:
        st.subheader("Executive KPI Snapshot")
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("Ending ARR", money(_metric(agent_results, "ending_arr")))
        col2.metric("Ending MRR", money(_metric(agent_results, "ending_mrr")))
        col3.metric("Ending Customers", number(_metric(agent_results, "ending_customers")))
        col4.metric("Funding Need", money(_metric(agent_results, "funding_need")))

        col5, col6, col7 = st.columns(3)
        col5.metric("Gross Margin", percent(_metric(agent_results, "gross_margin")))
        col6.metric("EBITDA Margin", percent(_metric(agent_results, "ebitda_margin")))
        col7.metric("LTV:CAC", f"{number(_metric(agent_results, 'ltv_cac'))}x")

        st.subheader("Generated Charts")
        for chart in agent_results.get("charts", []):
            st.image(chart, caption=Path(chart).name, use_container_width=True)

with tabs[3]:
    if not agent_results:
        st.info("Upload a workbook first.")
    else:
        st.subheader("Agent FP&A Analysis")
        finance = agent_results.get("finance_reasoning", {})
        for title, section in finance.items():
            if title == "recommended_actions":
                continue
            st.markdown(f"### {title.replace('_', ' ').title()}")
            for bucket in ["Facts", "Interpretation", "Recommendations"]:
                st.markdown(f"**{bucket}**")
                _render_bullets(section.get(bucket, []))

        st.markdown("### Recommended Actions")
        _render_bullets(finance.get("recommended_actions", []))

with tabs[4]:
    if not agent_results:
        st.info("Upload a workbook first.")
    else:
        st.subheader("Devil's Advocate Review")
        challenge = agent_results.get("devils_advocate", {})
        for label, key in [
            ("Investor Objections", "investor_objections"),
            ("Weak Assumptions", "weak_assumptions"),
            ("Data Gaps", "data_gaps"),
            ("Questions Management Must Answer", "questions_management_must_answer"),
        ]:
            st.markdown(f"### {label}")
            _render_bullets(challenge.get(key, []))

with tabs[5]:
    if not agent_results:
        st.info("Upload a workbook first.")
    else:
        report_path = Path(agent_results.get("reports", {}).get("markdown_report", ""))
        report = report_path.read_text(encoding="utf-8") if report_path.exists() else ""
        metrics = agent_results.get("metrics", {})
        pdf_bytes = _build_pdf(report, agent_results)
        pptx_bytes = _build_pptx(agent_results)

        st.subheader("Quality Review")
        quality = agent_results.get("quality_review", {})
        col1, col2 = st.columns([1, 3])
        col1.metric("Quality Score", f"{quality.get('quality_score', 0)}/100")
        with col2:
            _status_badge(quality.get("status", "Needs Review"))

        st.write("Checklist")
        for item in quality.get("checklist", []):
            st.write(f"- {'Pass' if item.get('passed') else 'Review'}: {item.get('item')}")

        st.write("Final review notes")
        _render_bullets(quality.get("final_review_notes", []))

        st.subheader("Download Outputs")
        st.download_button(
            "Download Metrics JSON",
            data=_json_text(metrics),
            file_name="cloud_ai_metrics.json",
            mime="application/json",
        )
        st.download_button(
            "Download Agent Results JSON",
            data=_json_text(agent_results),
            file_name="cloud_ai_agent_results.json",
            mime="application/json",
        )
        st.download_button(
            "Download Board Report Markdown",
            data=report,
            file_name="cloud_ai_board_report.md",
            mime="text/markdown",
        )
        st.download_button(
            "Download Board Report PDF",
            data=pdf_bytes,
            file_name="cloud_ai_board_report.pdf",
            mime="application/pdf",
        )
        st.download_button(
            "Download Board Deck PPTX",
            data=pptx_bytes,
            file_name="cloud_ai_board_deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        for chart in agent_results.get("charts", []):
            chart_path = Path(chart)
            if chart_path.exists():
                st.download_button(
                    f"Download {chart_path.name}",
                    data=chart_path.read_bytes(),
                    file_name=chart_path.name,
                    mime="image/png",
                )
